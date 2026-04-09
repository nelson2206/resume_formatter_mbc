from pptx import Presentation
import os

# Ruta absoluta con la nueva plantilla
TEMPLATE = r"C:\Users\nebernal\OneDrive - Indra\Documentos\CV_To_PPT_App\templates\CV template.pptx"

print(f"Archivo existe: {os.path.exists(TEMPLATE)}")

prs = Presentation(TEMPLATE)

for slide_idx, slide in enumerate(prs.slides):
    print(f"\n=== SLIDE {slide_idx} | Layout: {slide.slide_layout.name} | Shapes: {len(slide.shapes)} ===\n")
    for i, shape in enumerate(slide.shapes):
        print(f"--- Shape {i} | Name: \"{shape.name}\" | Type: {shape.shape_type} ---")
        if shape.has_text_frame:
            tf = shape.text_frame
            full = tf.text
            print(f"  Full text: {repr(full[:300])}")
            for pi, para in enumerate(tf.paragraphs):
                runs_info = []
                for ri, run in enumerate(para.runs):
                    font = run.font
                    try:
                        color = str(font.color.rgb) if font.color and font.color.type else 'inherit'
                    except:
                        color = 'inherit'
                    runs_info.append(f"run{ri}={repr(run.text[:40])}[bold={font.bold},sz={font.size},col={color}]")
                print(f"  Para {pi}: {repr(para.text[:100])} | {' | '.join(runs_info)}")
        else:
            print(f"  (sin text frame)")
        print()
