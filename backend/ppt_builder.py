"""
ppt_builder.py — Generador de PPT basado en la plantilla con marcadores explícitos.

Esta versión soporta plantillas que usen los siguientes marcadores de texto:
{{NombreApellido}}
{{Rol_Seniority}}
{{Formacion_academica}}
{{conocimientos_certificaciones_idiomas}}
{{certificaciones}}
{{Resumen}}
{{Responsabilidades}} (este se expandirá a múltiples bullets)

La principal ventaja es que preserva el formato exacto del texto marcado.
"""

import copy
import re
import math
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

# Espacio de nombres DrawingML y atributo xml:space para preservar espacios entre runs
_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_XML_SPACE = "{http://www.w3.org/XML/1998/namespace}space"


def _parse_segmentos_negrita(texto):
    """
    Divide un texto con marcadores Markdown **negrita** en una lista de
    tuplas (fragmento, es_negrita). Los asteriscos sueltos sin pareja se
    eliminan del texto normal para evitar artefactos visuales.

    Ej: "Lideró **SAP** en banca" -> [("Lideró ", False), ("SAP", True), (" en banca", False)]
    """
    if not texto:
        return [("", False)]

    partes = re.split(r"\*\*(.+?)\*\*", texto)
    segmentos = []
    for idx, parte in enumerate(partes):
        es_negrita = idx % 2 == 1
        if not es_negrita:
            parte = parte.replace("**", "")  # limpiar asteriscos huérfanos
        if parte == "":
            continue
        segmentos.append((parte, es_negrita))

    if not segmentos:
        return [("", False)]
    return segmentos

def generar_ppt_desde_plantilla(template_path: str, output_path: str, perfil: dict):
    """
    Genera un PPT llenando la plantilla con los datos del perfil (usando placeholders).
    """
    prs = Presentation(template_path)
    
    reemplazos_texto = {
        "{{Resumen}}": perfil.get("resumen_profesional", ""),
        "{{idiomas}}": perfil.get("idiomas", "")
    }

    reemplazos_lista = {
        "{{Responsabilidades}}": perfil.get("experiencia_profesional", []),
        "{{conocimientos_clave}}": perfil.get("conocimientos_clave", []),
        "{{certificaciones}}": perfil.get("certificaciones", []),
        "{{Formacion_academica}}": perfil.get("formacion_academica", [])
    }
    
    # Iteramos sobre todos los slides (normalmente será solo 1)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
                
            tf = shape.text_frame
            # ¿Es el cuadro de experiencia? Guardamos sus dimensiones para repartir
            # los bullets verticalmente y evitar el espacio en blanco cuando hay pocos.
            es_shape_experiencia = "{{Responsabilidades}}" in tf.text
            box_h_emu = shape.height
            box_w_emu = shape.width
            # Párrafos de experiencia creados (para el espaciado adaptativo)
            exp_paragraphs = []
            exp_textos = []
            # Lista de párrafos que necesitaremos borrar (porque los expandimos)
            paras_to_remove = []
            if not tf.paragraphs:
                continue
            
            for p_idx, para in enumerate(list(tf.paragraphs)):
                texto_parrafo = para.text
                
                # REGLAS DE FORMATO MINSAIT:
                # 1. Nombre (24) en COLOR BLANCO
                if "{{NombreApellido}}" in texto_parrafo:
                    pPr = copy.deepcopy(para._p.get_or_add_pPr())
                    para.clear()
                    para._p.append(pPr)
                    r1 = para.add_run()
                    r1.text = perfil.get("nombre", "")
                    r1.font.size = Pt(24)
                    r1.font.color.rgb = RGBColor(255, 255, 255)
                    continue
                
                # 2. Rol (18) en COLOR BLANCO
                if "{{Rol_Seniority}}" in texto_parrafo:
                    pPr = copy.deepcopy(para._p.get_or_add_pPr())
                    para.clear()
                    para._p.append(pPr)
                    r2 = para.add_run()
                    r2.text = perfil.get("rol_seniority", "")
                    r2.font.size = Pt(18)
                    r2.font.color.rgb = RGBColor(255, 255, 255)
                    continue
                    
                # 3. Resumen (12) — con resaltado en negrita
                if "{{Resumen}}" in texto_parrafo:
                    pPr = copy.deepcopy(para._p.get_or_add_pPr())
                    para.clear()
                    para._p.append(pPr)
                    for seg_text, es_negrita in _parse_segmentos_negrita(perfil.get("resumen_profesional", "")):
                        run = para.add_run()
                        run.text = seg_text
                        run.font.size = Pt(12)
                        run.font.bold = es_negrita
                        # Preservar espacios al inicio/fin del fragmento entre runs
                        t_elem = run._r.find(qn("a:t"))
                        if t_elem is not None:
                            t_elem.set(_XML_SPACE, "preserve")
                    continue
                
                # 4. Chequeo de reemplazos simples restantes (si hubiera)
                reemplazo_simple = False
                if any(t in texto_parrafo for t in reemplazos_texto.keys()):
                    pPr = copy.deepcopy(para._p.get_or_add_pPr())
                    nuevo_texto = texto_parrafo
                    for token, val in reemplazos_texto.items():
                        if token in nuevo_texto:
                            nuevo_texto = nuevo_texto.replace(token, str(val) if val else "")
                    para.clear()
                    para._p.append(pPr)
                    para.text = nuevo_texto
                    for r in para.runs:
                        r.font.size = Pt(12)
                    reemplazo_simple = True
                        
                if reemplazo_simple:
                    continue
                    
                # 2. Chequeo de expansión de listas / bullets
                for token, items_lista in reemplazos_lista.items():
                    if token in texto_parrafo:
                        p_elem = para._p
                        # Crear un nuevo párrafo por cada item. Si la lista está vacía
                        # (ej. {{certificaciones}} ya fusionado en formación), el placeholder
                        # se elimina sin dejar línea en blanco.
                        nuevos = []
                        for item_str in items_lista:
                            nuevo_p = _clonar_parrafo_con_texto(para, item_str)
                            p_elem.addprevious(nuevo_p)
                            nuevos.append(nuevo_p)

                        if token == "{{Responsabilidades}}":
                            exp_paragraphs = nuevos
                            exp_textos = list(items_lista)

                        # Marcar el placeholder original para borrarlo después
                        paras_to_remove.append(p_elem)
                        break
            
            # Limpiamos los párrafos originales que fueron expandidos
            for p_elem in paras_to_remove:
                parent = p_elem.getparent()
                if parent is not None:
                    parent.remove(p_elem)

            # Espaciado adaptativo: si es el cuadro de experiencia y quedó holgura,
            # repartimos el espacio en blanco como separación entre bullets.
            if es_shape_experiencia and exp_paragraphs:
                spc = _espaciado_adaptativo_pt(exp_textos, box_h_emu, box_w_emu)
                if spc > 0:
                    for nuevo_p in exp_paragraphs:
                        pPr = nuevo_p.find(f"{{{_NS_A}}}pPr")
                        if pPr is None:
                            pPr = etree.Element(f"{{{_NS_A}}}pPr")
                            nuevo_p.insert(0, pPr)
                        _set_space_before(pPr, spc)

    prs.save(output_path)


def _set_space_before(pPr, pt):
    """Inserta/actualiza <a:spcBef> en un pPr respetando el orden del esquema
    DrawingML (spcBef va tras lnSpc y antes de spcAft / viñetas)."""
    ns = _NS_A
    for el in pPr.findall(f"{{{ns}}}spcBef"):
        pPr.remove(el)
    spcBef = etree.Element(f"{{{ns}}}spcBef")
    spcPts = etree.SubElement(spcBef, f"{{{ns}}}spcPts")
    spcPts.set("val", str(int(round(pt * 100))))  # centésimas de punto
    lnSpc = pPr.find(f"{{{ns}}}lnSpc")
    if lnSpc is not None:
        lnSpc.addnext(spcBef)
    else:
        pPr.insert(0, spcBef)


def _espaciado_adaptativo_pt(textos, box_height_emu, box_width_emu, font_pt=12.0):
    """
    Calcula el espacio (en puntos) a colocar ANTES de cada bullet para distribuir
    la experiencia y reducir el espacio en blanco cuando hay pocos bullets, SIN
    inventar contenido. Es conservador: estima por exceso la altura del texto para
    no provocar desbordes, y limita el espaciado máximo.
    """
    EMU_PT = 12700.0
    n = len([t for t in textos if t and t.strip()])
    if n == 0 or not box_height_emu or not box_width_emu:
        return 0.0

    box_h = (box_height_emu / EMU_PT) * 0.90      # margen de seguridad interno
    box_w = box_width_emu / EMU_PT
    char_w = font_pt * 0.55                        # ancho de carácter (estimación amplia)
    chars_per_line = max(1.0, box_w / char_w)
    line_h = font_pt * 1.30                         # alto de línea aprox.

    text_h = 0.0
    for t in textos:
        if not (t and t.strip()):
            continue
        limpio = t.replace("**", "")               # los asteriscos no se renderizan
        lineas = max(1, math.ceil(len(limpio) / chars_per_line))
        text_h += lineas * line_h

    slack = box_h - text_h
    if slack <= 0:
        return 0.0

    spc = slack / n
    return max(0.0, min(spc, 26.0))                 # tope de 26pt para que no quede excesivo


def _clonar_parrafo_con_texto(para_referencia, texto: str) -> etree._Element:
    """
    Devuelve un nuevo elemento <a:p> copiado del formato del párrafo de referencia.
    El `texto` puede contener marcadores Markdown **negrita**, que se renderizan como
    runs independientes en negrita preservando el resto del formato Minsait.
    """
    ns = _NS_A
    p_new = etree.Element(f"{{{ns}}}p")

    # Obtener el elemento XML del párrafo de referencia
    p_ref_elem = para_referencia._p

    # 1. Copiar pPr (propiedades de párrafo - ej. viñetas, márgenes)
    pPr_ref = p_ref_elem.find(f"{{{ns}}}pPr")
    if pPr_ref is not None:
        pPr_nuevo = copy.deepcopy(pPr_ref)

        # AJUSTE DE IDENTACIÓN: Para que la segunda línea se alinee con el inicio de las letras
        # marL es el margen izquierdo del texto, indent es el desplazamiento de la primera línea (negativo para viñeta)
        # 457200 EMU = 12.7mm = 1/2 inch.
        pPr_nuevo.set("marL", "228600")
        pPr_nuevo.set("indent", "-228600")

        p_new.append(pPr_nuevo)

    # 2. Localizar el primer run para copiar su formato base
    runs_existentes = p_ref_elem.findall(f"{{{ns}}}r")
    rPr_base = None
    if runs_existentes:
        rPr_existente = runs_existentes[0].find(f"{{{ns}}}rPr")
        if rPr_existente is not None:
            rPr_base = copy.deepcopy(rPr_existente)

    if rPr_base is None:
        rPr_base = etree.Element(f"{{{ns}}}rPr")

    # Forzar el formato Minsait: 12pt (1200 centésimas de punto)
    rPr_base.set("sz", "1200")

    # 3. Crear un run por cada fragmento, aplicando negrita donde corresponda
    for seg_text, es_negrita in _parse_segmentos_negrita(texto):
        r_new = etree.SubElement(p_new, f"{{{ns}}}r")
        rPr_seg = copy.deepcopy(rPr_base)
        if es_negrita:
            rPr_seg.set("b", "1")
        else:
            rPr_seg.attrib.pop("b", None)
        r_new.append(rPr_seg)
        t_new = etree.SubElement(r_new, f"{{{ns}}}t")
        t_new.set(_XML_SPACE, "preserve")  # preservar espacios entre fragmentos
        t_new.text = seg_text

    return p_new
