import io
import pdfplumber
import re
import fitz  # PyMuPDF
try:
    import easyocr
    import numpy as np
except ImportError:
    easyocr = None
    np = None

from PIL import Image

# Marca que devuelve el extractor cuando un PDF no tiene texto util (escaneado)
# y no hay OCR local disponible. main.py la usa para enviar las paginas como
# imagen al modelo multimodal en lugar de descartar el CV.
MARCA_ESCANEADO = "[[PDF_ESCANEADO]]"

# Inicializar lector OCR (Español e Inglés) - Singleton para evitar recargas
_reader = None

def get_ocr_reader():
    global _reader
    if easyocr is None:
        print("[extractor] EasyOCR no está instalado. El soporte para PDFs escaneados está desactivado.")
        return None
        
    if _reader is None:
        print("[extractor] Inicializando EasyOCR (esto puede tardar la primera vez)...")
        try:
            _reader = easyocr.Reader(['es', 'en'], gpu=False) # Forzamos CPU por compatibilidad
        except Exception as e:
            print(f"[extractor] No se pudo inicializar EasyOCR: {e}")
            _reader = None
    return _reader


def extraer_texto(filename: str, file_bytes: bytes) -> str:
    """
    Detecta el tipo de archivo por extensión y extrae texto.
    Devuelve string limpio o cadena vacía si falla.
    """
    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""

    if ext == "pdf":
        return _extraer_pdf(file_bytes)
    elif ext in ("docx", "doc"):
        return _extraer_docx(file_bytes)
    elif ext in ("pptx", "ppt"):
        return _extraer_pptx(file_bytes)
    elif ext == "txt":
        return _extraer_txt(file_bytes)
    else:
        return ""


def _extraer_pdf(file_bytes: bytes) -> str:
    """Extrae texto de PDF usando un algoritmo de detección de columnas o OCR si es imagen."""
    texto_final = ""
    try:
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            total_chars = 0
            for page in pdf.pages:
                total_chars += len(page.chars)
            
            # Si hay muy pocos caracteres (ej: < 100), es probable que sea una imagen.
            if total_chars < 100:
                print(f"[extractor] Detectado PDF de imagen ({total_chars} chars). Iniciando OCR...")
                return _extraer_ocr(file_bytes)

            for page in pdf.pages:
                words = page.extract_words(x_tolerance=3, y_tolerance=3)
                if not words: continue

                height = float(page.height)
                width = float(page.width)
                header_y_limit = height * 0.25
                column_split_x = width * 0.40

                header_words = [w for w in words if w['bottom'] <= header_y_limit]
                body_words = [w for w in words if w['bottom'] > header_y_limit]

                left_words = [w for w in body_words if w['x1'] <= column_split_x]
                right_words = [w for w in body_words if w['x1'] > column_split_x]

                def reconstruct(word_list):
                    if not word_list: return ""
                    sorted_words = sorted(word_list, key=lambda w: (w['top'], w['x0']))
                    lines = []
                    if not sorted_words: return ""
                    
                    current_line = [sorted_words[0]['text']]
                    last_top = sorted_words[0]['top']
                    
                    for i in range(1, len(sorted_words)):
                        w = sorted_words[i]
                        if abs(w['top'] - last_top) < 3:
                            current_line.append(w['text'])
                        else:
                            lines.append(" ".join(current_line))
                            current_line = [w['text']]
                            last_top = w['top']
                    lines.append(" ".join(current_line))
                    return "\n".join(lines)

                texto_final += "--- SECCION ENCABEZADO ---\n" + reconstruct(header_words) + "\n\n"
                texto_final += "--- SECCION LATERAL (DATOS/CONOCIMIENTOS) ---\n" + reconstruct(left_words) + "\n\n"
                texto_final += "--- SECCION PRINCIPAL (EXPERIENCIA) ---\n" + reconstruct(right_words) + "\n\n"
                texto_final += "\n"
    except Exception as e:
        print(f"[extractor] Error leyendo PDF con columnas: {e}. Reintentando con OCR...")
        return _extraer_ocr(file_bytes)
    
    # Texto ilegible (caracteres de reemplazo por mala codificacion) o practicamente
    # vacio: se lee el PDF como imagen. El literal de este check estaba vacio, con lo
    # que cualquier PDF de menos de 500 caracteres caia aqui.
    if "�" in texto_final or len(texto_final.strip()) < 200:
        return _extraer_ocr(file_bytes)

    return _limpiar_texto(texto_final)


def _extraer_ocr(file_bytes: bytes) -> str:
    """Convierte PDF en imágenes y extrae texto con OCR."""
    texto_ocr = ""
    try:
        reader = get_ocr_reader()
        if reader is None:
            return MARCA_ESCANEADO
            
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2)) # Zoom 2x para mejor OCR
            img_data = pix.tobytes("png")
            
            # Convertir a formato que EasyOCR entiende (numpy array)
            img = Image.open(io.BytesIO(img_data))
            img_np = np.array(img)
            
            # Leer texto
            result = reader.readtext(img_np, detail=0, paragraph=True)
            texto_ocr += f"\n--- Pagina {page_num + 1} (OCR) ---\n"
            texto_ocr += "\n".join(result) + "\n"
            
        doc.close()
        return _limpiar_texto(texto_ocr)
    except Exception as e:
        print(f"[extractor] Error crítico en OCR: {e}")
        return ""


def _extraer_pptx(file_bytes: bytes) -> str:
    """Extrae texto de diapositivas de PowerPoint (.pptx)."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        texto_slides = []
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                texto_slides.append("\n".join(slide_text))
        return _limpiar_texto("\n\n--- Slide ---\n\n".join(texto_slides))
    except Exception as e:
        print(f"[extractor] Error leyendo PPTX: {e}")
        return ""


def _extraer_docx(file_bytes: bytes) -> str:
    """Extrae texto de archivos .docx usando python-docx."""
    try:
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        parrafos = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        parrafos.append(cell_text)
        return _limpiar_texto("\n".join(parrafos))
    except Exception as e:
        print(f"[extractor] Error leyendo DOCX: {e}")
        return ""


def _extraer_txt(file_bytes: bytes) -> str:
    """Extrae texto de archivos planos .txt."""
    try:
        texto = file_bytes.decode("utf-8", errors="ignore")
        return _limpiar_texto(texto)
    except Exception as e:
        print(f"[extractor] Error leyendo TXT: {e}")
        return ""


def _limpiar_texto(texto: str) -> str:
    """Normaliza espacios, elimina caracteres de control innecesarios."""
    texto = texto.replace("\t", " ")
    texto = re.sub(r" {2,}", " ", texto)
    texto = re.sub(r"\n{3,}", "\n\n", texto)
    return texto.strip()


def extraer_imagenes_pdf(file_bytes: bytes, max_paginas: int = 4, ancho_max: int = 1600) -> list:
    """
    Renderiza las paginas de un PDF a JPEG para enviarlas a un modelo multimodal.

    Se usa con CVs escaneados, donde no hay texto que extraer. Limita paginas y
    ancho para no disparar el coste ni el tamano de la peticion.

    Returns: lista de bytes JPEG, una por pagina. Vacia si el PDF no se puede abrir.
    """
    imagenes = []
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        for page_num in range(min(len(doc), max_paginas)):
            page = doc.load_page(page_num)
            ancho_pt = float(page.rect.width) or 1.0
            zoom = min(2.0, ancho_max / ancho_pt)
            pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom))
            img = Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGB")
            buf = io.BytesIO()
            img.save(buf, format="JPEG", quality=80, optimize=True)
            imagenes.append(buf.getvalue())
        doc.close()
        print("[extractor] PDF escaneado: %d pagina(s) renderizadas para lectura visual." % len(imagenes))
    except Exception as e:
        print("[extractor] No se pudieron renderizar las paginas del PDF: %s" % e)
        return []
    return imagenes
